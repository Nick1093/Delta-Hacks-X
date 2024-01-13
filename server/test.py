from pptx import Presentation


def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    slides_content = []  # Array to hold the content of each slide

    for slide in prs.slides:
        slide_text = ""  # String to hold content of each slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"  # Append text from each shape to the slide text

        slides_content.append(slide_text.strip())  # Add the slide's content to the array

    return slides_content

# Call the function and pass the PowerPoint file name
slides_content = extract_text_from_pptx('Lecture 6.pptx')

# Print the content of each slide
for slide_number, slide_content in enumerate(slides_content, 1):
    print(f"Content of Slide {slide_number}:\n{slide_content}\n")
