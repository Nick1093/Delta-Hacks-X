from fastapi import FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect
from pptx import Presentation
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List
import coTest
import cohere
import asyncio
import json
# from WebSocket import ConnectionManager
# from coTest import *

app = FastAPI()

# COHERE_API_KEY = os.environ["COHERE_API_KEY"]
# co = cohere.Client(COHERE_API_KEY)

# Set up CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)

class ConnectionManager:
    def __init__(self):
        self.active_connections = []
    
    async def connect(self, websocket: WebSocket):
        await websocket.accept()
        self.active_connections.append(websocket)

    async def send_personal_message(self, message: str, websocket: WebSocket):
        await websocket.send_text(message)
    
    def disconnect(self, websocket: WebSocket):
        self.active_connections.remove(websocket)
    
    # data received is an array ["element1", "element2"] where each element is a slide's worth of text
    # this function should return 
    async def generateContent(self, pptx_data):
        return coTest.generate_reel_content(pptx_data)
    
    async def generate_and_send(self, scraped_content, websocket):
        first_run = scraped_content["slides_content"][:len(scraped_content["slides_content"]) // 2]
        new_reel_content = await self.generateContent(first_run)
        await websocket.send_json({"title": scraped_content["title"], "reel_content": new_reel_content})

manager = ConnectionManager()

@app.get("/")
def read_root():
    return {"Hello": "World"}


@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    await manager.connect(websocket)
    try:
        while True:
            data = await websocket.receive_text()
            reels_content = json.loads(data)
            for scraped_content in reels_content:
                await manager.generate_and_send(scraped_content, websocket)
    except WebSocketDisconnect:
        print("Error occurred, disconnecting...")
        manager.disconnect(websocket)
        await manager.broadcast("A client disconnected.")



# helper method to extract text from each slide for the pptx
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    ppt_title = prs.core_properties.title
    slides_content = []  # Array to hold the content of each slide
    slide_data = {"title": ppt_title, "slides_content":[]}

    try:
        for slide in prs.slides:
            slide_text = ""  # String to hold content of each slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"  # Append text from each shape to the slide text
                slide_text = slide_text.strip()

            if slide_text:  # Check if the slide text is not empty
                slides_content.append(slide_text)  # Add the slide's content to the array
            
        slide_data["slides_content"] = slides_content
    except:
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slides_content.append(run.text)
        except:
            print("Content retrieval unsuccessful")

        slides_content.append(slide_text.strip())  # Add the slide's content to the array
        slide_data["slides_content"] = slides_content

    return slide_data

@app.post("/upload_pptx/")
async def upload_pptx(uploaded_files: List[UploadFile] = File(...)):
    try:
        # content of the reels
        reels_content = []

        for i in range(len(uploaded_files)):
            with open("temp.pptx", "wb") as f:
                f.write(uploaded_files[i].file.read())

            # Extract text from the PowerPoint file
            powerpoint_content = extract_text_from_pptx("temp.pptx")

            reels_content.append(powerpoint_content)

        #Return the extracted text as a response
        return {"success": True, "slides_content": reels_content}
    except Exception as e:
        return {"success": False, "error_message": str(e)}
    
# @app.get("/generate_response/")
# async def generate_response(chat_history: List[map], new_prompt: str):
#     try:
#         response = co.chat(
#             message=new_prompt, 
#             chat_history=chat_history,
#             model="command", 
#             temperature=0.9
#         )

#         answer = response.text
        
#         return {"success": True, "response": answer}
#     except Exception as e:
#         return {"success": False, "error_message": str(e)}


# content of the reels
            # reels_content = []
            # uploaded_files = await websocket.receive_json()
            # for i in range(len(uploaded_files)):
            #     with open("temp.pptx", "wb") as f:
            #         f.write(uploaded_files[i].file.read())

            #     # Extract text from the PowerPoint file
            #     slides_content = extract_text_from_pptx("temp.pptx")

            #     reels_content.append(slides_content)
    
    # tasks = [asyncio.create_task(process_data(item, websocket)) for item in data]
# websocket
# @app.websocket("/ws")
# async def websocket_endpoint(websocket: WebSocket):
#     await manager.connect(websocket)

#     try:
#         while True:
#             data = await websocket.receive_text()
#             # Parse JSON string back into an array
#             reels_content = json.loads(data)
#             for scraped_content in reels_content:
#                 first_run = scraped_content["slides_content"][:len(scraped_content["slides_content"])//2]
#                 new_reel_content_1 = await manager.generateContent(first_run)
#                 print("Content Generated: ", new_reel_content_1)
#                 await websocket.send_json({"title": scraped_content["title"], "reel_content": new_reel_content_1})
#     except WebSocketDisconnect:
#         print("Error occured, disconecting...")
#         manager.disconnect(websocket)
#         await manager.broadcast("A client disconnected.")